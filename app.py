from flask import Flask, request, jsonify, send_file
import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import os

app = Flask(__name__)

@app.route("/")
def home():
    return "Document API is running!"

@app.route("/create-excel", methods=["POST"])
def create_excel():
    data = request.get_json()
    title = data.get("title", "output")
    rows = data.get("rows")

    if not rows:
        return jsonify({"error": "データが空です"}), 400

    try:
        df = pd.DataFrame(rows[1:], columns=rows[0])
        filename = f"{title}.xlsx"
        df.to_excel(filename, index=False)
        return send_file(filename, as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/create-word", methods=["POST"])
def create_word():
    data = request.get_json()
    title = data.get("title", "output")
    paragraphs = data.get("paragraphs", [])

    try:
        doc = Document()
        for para in paragraphs:
            doc.add_paragraph(para)

        filename = f"{title}.docx"
        doc.save(filename)
        return send_file(filename, as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/create-pptx", methods=["POST"])
def create_pptx():
    data = request.get_json()
    title = data.get("title", "output")
    slides = data.get("slides", [])

    try:
        prs = Presentation()
        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_placeholder = slide.shapes.title
            content_placeholder = slide.placeholders[1]

            title_placeholder.text = slide_data.get("title", "")
            content_placeholder.text = slide_data.get("content", "")

        filename = f"{title}.pptx"
        prs.save(filename)
        return send_file(filename, as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port)
